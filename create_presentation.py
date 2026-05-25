#!/usr/bin/env python3
"""Generate IR-Sensor project takeover presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)  # Deep navy
ACCENT       = RGBColor(0xE9, 0x45, 0x60)  # Red accent
ACCENT2      = RGBColor(0x0F, 0xC3, 0xFF)  # Cyan accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY     = RGBColor(0x99, 0x99, 0x99)
DARK_CARD    = RGBColor(0x22, 0x22, 0x3A)
GREEN        = RGBColor(0x2E, 0xCC, 0x71)
YELLOW       = RGBColor(0xF1, 0xC4, 0x0F)
ORANGE       = RGBColor(0xE6, 0x7E, 0x22)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, bg_color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(bullet_color[0], bullet_color[1], bullet_color[2]))
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_numbered_item(slide, left, top, width, height, number, title, description,
                      num_color=ACCENT, title_size=20, desc_size=15):
    add_textbox(slide, left, top, Inches(0.6), Inches(0.5),
                f"{number}.", font_size=title_size, color=num_color, bold=True)
    add_textbox(slide, left + Inches(0.7), top, width - Inches(0.7), Inches(0.4),
                title, font_size=title_size, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.7), top + Inches(0.45), width - Inches(0.7), height - Inches(0.45),
                description, font_size=desc_size, color=LIGHT_GRAY)

# ── Create Presentation ────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(2.8), Inches(11.3), ACCENT)
add_textbox(slide, Inches(1), Inches(3), Inches(11.3), Inches(1.2),
            "IR-Sensor Projektübernahme", font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.8),
            "Status, Herausforderungen & Vorgehensplan", font_size=24, color=LIGHT_GRAY)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
            "Präsentation an die Projektleitung", font_size=16, color=MID_GRAY)
add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
            "2025", font_size=14, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "Agenda", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(2), ACCENT)

agenda_items = [
    ("1", "Projekt & aktueller Stand"),
    ("2", "Identifizierte Kernprobleme"),
    ("3", "Priorisierte Arbeitsreihenfolge"),
    ("4", "AI-gestützte Entwicklungsmethodik"),
    ("5", "Infrastruktur: Jira & Virtual Machine"),
    ("6", "Team-Unterstützung & Knowledge"),
    ("7", "Messaufbau & Feldtests"),
    ("8", "Nächste Schritte & Projektplan"),
]
for i, (num, title) in enumerate(agenda_items):
    y = Inches(1.8) + i * Inches(0.65)
    add_card(slide, Inches(1), y, Inches(10.5), Inches(0.55))
    add_textbox(slide, Inches(1.3), y + Inches(0.05), Inches(0.5), Inches(0.45),
                num, font_size=20, color=ACCENT, bold=True)
    add_textbox(slide, Inches(1.9), y + Inches(0.05), Inches(9), Inches(0.45),
                title, font_size=20, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: Projekt & aktueller Stand
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "Projekt IR-Sensor — Aktueller Stand", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT)

# Left column - Status
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
            "✅ Bereits abgeschlossen", font_size=22, color=GREEN, bold=True)
status_items = [
    "Stakeholder-Anforderungen erfasst (vor Monaten)",
    "KI-gestützte Transformation in Systemanforderungen",
    "Verlinkung der Anforderungen erfolgt",
    "Parallele Vorarbeit an dem Projekt in den letzten Wochen",
    "4 Kernprobleme identifiziert und dokumentiert",
]
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4),
                status_items, font_size=16, color=LIGHT_GRAY)

# Right column - Context
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.5),
            "📋 Projektkontext", font_size=22, color=ACCENT2, bold=True)
context_items = [
    "Einzelkämpfer-Projekt (kein Team)",
    "Organisationsstrukturen müssen etabliert werden",
    "AI-Entwicklungsmethoden (AI Codex) als Enabler",
    "Ziel: Zuverlässiger IR-Sensor für Außenfeld-Messungen",
    "Hardware: EPC611, STM32L4, IR-LED, LIN-BUS",
]
add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(4),
                context_items, font_size=16, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: Kernprobleme Overview
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Identifizierte Kernprobleme", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT)

problems = [
    ("1", "EPC611-Abstürze", "Ca. alle 25 Min. → kein Signal mehr",
     "Workaround: Spannung weg, Neustart. Root-Cause muss gefunden werden.",
     ACCENT),
    ("2", "Teilverschattete Flächen", "Sensor-Übersteuerung in Matrix-Bereichen",
     "Keine Messung mehr möglich. IR-LED im Labor reproduzierbar, im Freifeld erwartet.",
     ORANGE),
    ("3", "LIN-BUS / Nachbarschaft", "UART-Zuordnung nach STM32F1→L4 Migration defekt",
     "MAX Richter löst parallel. Zeitrahmen: offen.",
     YELLOW),
    ("4", "Fehlerlogging & EEPROM", "Kein persistenter Fehler-Speicher, kein ordentliches Logging",
     "EEPROM + CAN-Schnittstelle müssen genutzt werden.",
     GREEN),
]

for i, (num, title, subtitle, desc, color) in enumerate(problems):
    y = Inches(1.8) + i * Inches(1.35)
    add_card(slide, Inches(0.8), y, Inches(11.7), Inches(1.2))
    # Number circle
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
                num, font_size=28, color=color, bold=True)
    # Title
    add_textbox(slide, Inches(1.6), y + Inches(0.05), Inches(9), Inches(0.4),
                title, font_size=22, color=WHITE, bold=True)
    # Subtitle
    add_textbox(slide, Inches(1.6), y + Inches(0.4), Inches(9), Inches(0.35),
                subtitle, font_size=16, color=color)
    # Description
    add_textbox(slide, Inches(1.6), y + Inches(0.75), Inches(9), Inches(0.4),
                desc, font_size=14, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: Problem 1 - EPC611
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Problem 1: EPC611-Abstürze (Prio 🔴)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT)

add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
            "Symptom", font_size=22, color=ACCENT, bold=True)
symptom_items = [
    "Absturz ca. alle 25 Minuten",
    "Nach Absturz: keine Kommunikation mehr",
    "Bisheriger Fix: Spannung trennen → Neustart",
    "Zeigt auf tiefer liegende, systematische Probleme",
]
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4),
                symptom_items, font_size=16, color=LIGHT_GRAY)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.5),
            "Erforderliche Maßnahmen", font_size=22, color=GREEN, bold=True)
action_items = [
    "Root-Cause-Analyse (Memory Leak? Watchdog? Hardware?)",
    "Logging-Infrastruktur als Voraussetzung etablieren",
    "Stabile Basis für alle weiteren Messmodi",
    "Zuverlässige Lösung statt Workaround nötig",
]
add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(4),
                action_items, font_size=16, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: Problem 2 - Verschattung
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Problem 2: Teilverschattete Flächen (Prio 🟠)", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ORANGE)

add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
            "Problem", font_size=22, color=ORANGE, bold=True)
prob_items = [
    "Bei teilverschatteten Bodenflächen",
    "Sensor-Matrix wird in Bereichen übersteuert",
    "Folge: Keine Messung mehr in diesen Bereichen",
    "IR-LED im Labor: Problem nicht zuverlässig reproduzierbar",
    "Annahme: Im Freifeld (Sommer) reproduzierbar",
]
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4),
                prob_items, font_size=16, color=LIGHT_GRAY)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.5),
            "Lösungsansatz", font_size=22, color=GREEN, bold=True)
sol_items = [
    "Messaufbau im Außenfeld Kamen erforderlich",
    "Sobald erste Messdaten empfangen werden können",
    "Unterstützung durch Tobias Lippe, Lukas Igges, Lukas Dreier (BaLi)",
    "Adaptive Sensor-Kalibrierung bei Mischlicht",
]
add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(4),
                sol_items, font_size=16, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: Problem 3 & 4
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Problem 3: LIN-BUS & Problem 4: Fehlerlogging", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), YELLOW)

# Problem 3
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5))
add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.4),
            "⚡ LIN-BUS / Nachbarschaftsproblem", font_size=20, color=YELLOW, bold=True)
lin_items = [
    "UART-Zuordnung nach STM32F1 → STM32L4 Migration passt nicht mehr",
    "LIN-BUS nur teilweise verfügbar",
    "MAX Richter löst parallel (Zeitrahmen: offen)",
]
add_bullet_list(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(2),
                lin_items, font_size=15, color=LIGHT_GRAY)

# Problem 4
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.5))
add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5), Inches(0.4),
            "📝 Fehlerlogging & EEPROM/CAN", font_size=20, color=GREEN, bold=True)
log_items = [
    "Module müssen Fehler persistent speichern können",
    "EEPROM-Nutzung für Fehlerprotokolle etablieren",
    "CAN-Schnittstelle für Logging-Export nutzen",
    "Voraussetzung für alle Debugging-Aktivitäten",
]
add_bullet_list(slide, Inches(7.1), Inches(2.4), Inches(5.2), Inches(2),
                log_items, font_size=15, color=LIGHT_GRAY)

# Priority note
add_card(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5))
add_textbox(slide, Inches(1.1), Inches(4.8), Inches(11), Inches(0.4),
            "🎯 Arbeitsreihenfolge (Priorisierung)", font_size=22, color=WHITE, bold=True)
order_items = [
    "1. Logging-Infrastruktur (EEPROM + CAN) → Fundament für alle weiteren Arbeiten",
    "2. EPC611-Abstürze beheben (erst wenn Logging läuft → Root-Cause analysierbar)",
    "3. Messmodi TIM, Greyscale, UDN → erst wenn Basis stabil",
    "4. Verschattungsproblem → erst wenn Messdaten fließen (Außenfeld Kamen)",
    "5. LIN-BUS → parallel durch MAX Richter",
]
add_bullet_list(slide, Inches(1.1), Inches(5.3), Inches(11), Inches(2),
                order_items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT2)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: AI-Entwicklungsmethodik
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "AI-Gestützte Entwicklungsmethodik", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT2)

add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
            "🤖 AI Codex & AI-Methoden", font_size=22, color=ACCENT2, bold=True)
ai_items = [
    "AI Codex als primäres Entwicklungstool",
    "Automatisierte Code-Generierung & Review",
    "KI-gestützte Anforderungsanalyse (bereits im Einsatz)",
    "Reduzierung manueller Arbeitsschritte",
    "Schnellere Iteration bei Problembehebung",
]
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4),
                ai_items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT2)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.5),
            "📊 Vorteile der AI-Methodik", font_size=22, color=GREEN, bold=True)
adv_items = [
    "Ein-Personen-Projekt effizient bearbeiten",
    "Automatisierte Issue-Erstellung & Tracking",
    "Time-Tracking für Zeitablauf-Analysen",
    "Strukturierte Projektarbeit ohne Team-Overhead",
    "Reproduzierbare Workflows & Dokumentation",
]
add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(4),
                adv_items, font_size=16, color=LIGHT_GRAY, bullet_color=GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Infrastruktur
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Infrastruktur: Jira Board & Virtual Machine", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT2)

# Jira
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.6))
add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.4),
            "📋 Jira IRV4 Board", font_size=22, color=ACCENT, bold=True)
jira_items = [
    "Separates Board für IR-Sensor (IRV4)",
    "Automatischer Abruf von Issues & Stories",
    "Time-Tracking für Zeitablauf-Dokumentation",
    "In wenigen Tagen: Aussagen zum Zeitverbrauch",
    "Integration mit AI-Entwicklungspipeline",
]
add_bullet_list(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(2),
                jira_items, font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT)

# VM
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.6))
add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5), Inches(0.4),
            "🖥️ Virtual Machine", font_size=22, color=ACCENT2, bold=True)
vm_items = [
    "Ubuntu 26.04 Server",
    "4 CPUs / 32 GB RAM / ~80 GB (dynamisch)",
    "Automatische Jira-Ticket-Vorbereitung",
    "Analog zum BorgOS-Setup",
    "Auf einem der vorhandenen Server",
]
add_bullet_list(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(2),
                vm_items, font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT2)

# Purpose
add_card(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3))
add_textbox(slide, Inches(1.1), Inches(4.9), Inches(11), Inches(0.4),
            "🎯 Ziel der Infrastruktur", font_size=22, color=WHITE, bold=True)
infra_items = [
    "Nicht mehr alles händisch verwalten → Automatisierung durch AI-Pipeline",
    "Jira als Single Source of Truth für Tasks, Fortschritt & Zeit",
    "VM als Automatisierungs-Engine für Ticket-Preparation & BorgOS-Integration",
    "Sichtbarkeit für die Projektleitung: Echtzeit-Status & Time-Tracking",
]
add_bullet_list(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(1.5),
                infra_items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT2)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: Team & Knowledge
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Team-Unterstützung & Knowledge", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT2)

# Knowledge experts
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5))
add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.4),
            "🧠 Knowledge-Experten", font_size=22, color=YELLOW, bold=True)
add_textbox(slide, Inches(1.1), Inches(2.45), Inches(5), Inches(0.35),
            "Norbert", font_size=20, color=WHITE, bold=True)
add_textbox(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(0.35),
            "Datenblätter analysieren & zerlegen", font_size=15, color=LIGHT_GRAY)
add_textbox(slide, Inches(1.1), Inches(3.35), Inches(5), Inches(0.35),
            "Kevin", font_size=20, color=WHITE, bold=True)
add_textbox(slide, Inches(1.1), Inches(3.7), Inches(5), Inches(0.35),
            "Anforderungs-Review & Design-Review", font_size=15, color=LIGHT_GRAY)

# Field support
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.5))
add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5), Inches(0.4),
            "🔧 Feldunterstützung (BaLi)", font_size=22, color=GREEN, bold=True)
field_items = [
    "Tobias Lippe",
    "Lukas Igges",
    "Lukas Dreier",
    "Unterstützung bei Gelegenheit",
]
add_bullet_list(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(2),
                field_items, font_size=16, color=LIGHT_GRAY, bullet_color=GREEN)

# Parallel work
add_card(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3))
add_textbox(slide, Inches(1.1), Inches(4.9), Inches(11), Inches(0.4),
            "🔄 Parallelarbeit", font_size=22, color=WHITE, bold=True)
par_items = [
    "MAX Richter: LIN-BUS Problem (parallel, eigenständig)",
    "Norbert & Kevin: On-demand Knowledge (Datenblätter, Reviews)",
    "BaLi-Team: Messaufbau-Unterstützung bei Gelegenheit",
    "Hauptverantwortung & Entwicklung: Ich (allein) → Organisation ist Schlüssel",
]
add_bullet_list(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(1.5),
                par_items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT2)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: Messaufbau Kamen
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Messaufbau Kamen — Außenfeld", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT2)

add_card(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(11), Inches(0.5),
            "📍 Kontext", font_size=22, color=ACCENT2, bold=True)
kamen_items = [
    "Verschattungsproblem ist im Labor mit IR-LED nicht zuverlässig reproduzierbar",
    "Annahme: Im Freifeld (Sommer) lässt sich das Problem zuverlässig nachstellen",
    "Messaufbau in Kamen erforderlich, sobald erste Messdaten empfangen werden können",
    "Abhängigkeit: Messmodi TIM, Greyscale, UDN müssen zuerst funktionieren",
    "Unterstützung vor Ort durch BaLi-Team (Lippe, Igges, Dreier) bei Gelegenheit",
]
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(11), Inches(4),
                kamen_items, font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT2)

# Timeline note
add_card(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
add_textbox(slide, Inches(1.1), Inches(5.7), Inches(11), Inches(0.4),
            "⏱️ Voraussetzungen für den Messaufbau", font_size=20, color=YELLOW, bold=True)
pre_items = [
    "✓ Logging-Infrastruktur stabil  →  ✓ EPC611-Abstürze behoben  →  ✓ Messmodi funktionierend  →  📍 Außenfeld-Messung Kamen",
]
add_bullet_list(slide, Inches(1.1), Inches(6.2), Inches(11), Inches(0.7),
                pre_items, font_size=15, color=LIGHT_GRAY, bullet_color=YELLOW)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: Nächste Schritte
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Nächste Schritte & Projektplan", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3), ACCENT)

# Immediate
add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0))
add_textbox(slide, Inches(1.0), Inches(2.0), Inches(3.2), Inches(0.4),
            "Sofort", font_size=22, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.0), Inches(2.45), Inches(1.5), ACCENT)
imm_items = [
    "Jira IRV4 Board einrichten",
    "VM (Ubuntu 26.04) provisionieren",
    "AI-Pipeline konfigurieren",
    "Design-Phase starten",
    "Logging-Infrastruktur aufsetzen",
]
add_bullet_list(slide, Inches(1.0), Inches(2.7), Inches(3.2), Inches(4),
                imm_items, font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)

# Short term
add_card(slide, Inches(4.7), Inches(1.8), Inches(3.6), Inches(5.0))
add_textbox(slide, Inches(4.9), Inches(2.0), Inches(3.2), Inches(0.4),
            "Kurzfristig", font_size=22, color=ORANGE, bold=True)
add_accent_line(slide, Inches(4.9), Inches(2.45), Inches(1.5), ORANGE)
short_items = [
    "EPC611-Root-Cause analysieren",
    "EPC611-Absturz beheben",
    "Messmodi TIM/Greyscale/UDN",
    "EEPROM + CAN Logging operational",
    "Design-Review mit Kevin",
]
add_bullet_list(slide, Inches(4.9), Inches(2.7), Inches(3.2), Inches(4),
                short_items, font_size=14, color=LIGHT_GRAY, bullet_color=ORANGE)

# Medium term
add_card(slide, Inches(8.6), Inches(1.8), Inches(3.9), Inches(5.0))
add_textbox(slide, Inches(8.8), Inches(2.0), Inches(3.5), Inches(0.4),
            "Mittelfristig", font_size=22, color=GREEN, bold=True)
add_accent_line(slide, Inches(8.8), Inches(2.45), Inches(1.5), GREEN)
med_items = [
    "Messaufbau Kamen (Außenfeld)",
    "Verschattungsproblem angehen",
    "LIN-BUS (MAX Richter parallel)",
    "Feldtests im Sommer",
    "Projektplan finalisieren",
]
add_bullet_list(slide, Inches(8.8), Inches(2.7), Inches(3.5), Inches(4),
                med_items, font_size=14, color=LIGHT_GRAY, bullet_color=GREEN)

# Bottom note
add_card(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3))
add_textbox(slide, Inches(1.1), Inches(5.7), Inches(11), Inches(0.4),
            "📌 Projektplan folgt", font_size=22, color=WHITE, bold=True)
add_textbox(slide, Inches(1.1), Inches(6.15), Inches(11), Inches(0.5),
            "Detaillierter Projektplan mit Meilensteinen wird erstellt, sobald ich mich intensiver mit den einzelnen Aufgaben beschäftigen konnte.",
            font_size=16, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: Zusammenfassung / Danke
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(2.5), Inches(11.3), ACCENT)
add_textbox(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1),
            "Zusammenfassung", font_size=40, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

summary_items = [
    "4 Kernprobleme identifiziert & priorisiert",
    "Klare Arbeitsreihenfolge: Logging → Stabilität → Messmodi → Feldtests",
    "AI-gestützte Entwicklung für effiziente Ein-Personen-Arbeit",
    "Infrastruktur: Jira IRV4 + VM für Automatisierung & Tracking",
    "Knowledge-Experten & Feldunterstützung verfügbar",
    "Projektplan folgt nach intensiver Aufgabenanalyse",
]
add_bullet_list(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(3),
                summary_items, font_size=20, color=LIGHT_GRAY,
                bullet_color=ACCENT2, spacing=Pt(12))

add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
            "Fragen?", font_size=36, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────
output_path = "/home/bernd/Workbench/borg/IR-Sensor_Praesentation.pptx"
prs.save(output_path)
print(f"✅ Präsentation gespeichert: {output_path}")
